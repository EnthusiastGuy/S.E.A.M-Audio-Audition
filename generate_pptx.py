#!/usr/bin/env python3
"""
Generate a high-impact PPTX presentation for S.E.A.M Audio Audition.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x0D, 0x12)
BG_CARD      = RGBColor(0x16, 0x16, 0x20)
ACCENT       = RGBColor(0x00, 0xCC, 0xFF)   # Cyan
ACCENT2      = RGBColor(0x7C, 0x4D, 0xFF)   # Purple
ACCENT3      = RGBColor(0xFF, 0x6B, 0x6B)   # Coral
ACCENT_GREEN = RGBColor(0x6B, 0xCB, 0x77)   # Green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x44)
GOLD         = RGBColor(0xFF, 0xD9, 0x3D)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ───────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI', space_after=Pt(6)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_icon_circle(slide, left, top, size, color, text, font_size=20):
    shape = add_shape(slide, left, top, size, size, fill_color=color, shape_type=MSO_SHAPE.OVAL)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_feature_card(slide, left, top, width, height, icon_text, title, description, accent_color=ACCENT):
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    # Accent top bar
    add_shape(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Pt(3), fill_color=accent_color, shape_type=MSO_SHAPE.RECTANGLE)
    # Icon
    add_icon_circle(slide, left + Inches(0.3), top + Inches(0.35), Inches(0.55), accent_color, icon_text, font_size=16)
    # Title
    add_text_box(slide, left + Inches(1.0), top + Inches(0.3), width - Inches(1.3), Inches(0.4), title, font_size=16, color=WHITE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.3), top + Inches(0.85), width - Inches(0.6), height - Inches(1.0), description, font_size=11, color=LIGHT_GRAY)
    return card

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT):
    tf = add_rich_text_box(slide, left, top, width, height)
    for i, item in enumerate(items):
        p = add_paragraph(tf, f"  {item}", font_size=font_size, color=color, space_after=Pt(8))
        # Add a run at beginning for bullet
        from pptx.oxml.ns import qn
        r = p._p.r_lst[0] if p._p.r_lst else None
        # We'll prepend a colored bullet
    return tf

def add_comparison_row(slide, left, top, width, feature, seam_val, comp1_val, comp2_val, comp3_val, row_bg=None):
    col_w = width / 5
    h = Inches(0.4)
    if row_bg:
        add_shape(slide, left, top, width, h, fill_color=row_bg, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, left + Inches(0.1), top, col_w - Inches(0.1), h, feature, font_size=11, color=LIGHT_GRAY, bold=False)
    add_text_box(slide, left + col_w, top, col_w, h, seam_val, font_size=11, color=ACCENT_GREEN if seam_val == "✓" else (ACCENT3 if seam_val == "✗" else WHITE), bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + col_w*2, top, col_w, h, comp1_val, font_size=11, color=ACCENT_GREEN if comp1_val == "✓" else (ACCENT3 if comp1_val == "✗" else WHITE), bold=True if comp1_val in ("✓","✗") else False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + col_w*3, top, col_w, h, comp2_val, font_size=11, color=ACCENT_GREEN if comp2_val == "✓" else (ACCENT3 if comp2_val == "✗" else WHITE), bold=True if comp2_val in ("✓","✗") else False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + col_w*4, top, col_w, h, comp3_val, font_size=11, color=ACCENT_GREEN if comp3_val == "✓" else (ACCENT3 if comp3_val == "✗" else WHITE), bold=True if comp3_val in ("✓","✗") else False, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative gradient bars
for i, c in enumerate([ACCENT2, ACCENT, ACCENT_GREEN, GOLD]):
    add_shape(slide, Inches(1.5 + i * 2.7), Inches(0.15), Inches(2.2), Pt(4), fill_color=c, shape_type=MSO_SHAPE.RECTANGLE)

# Main title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
    "S.E.A.M", font_size=72, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER, font_name='Segoe UI Light')

add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
    "Segmented Evaluation & Audition Module", font_size=24, color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.0), Inches(3.7), Inches(3.333), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
    "The browser-based audition player for dynamic,\nbranching & loopable music", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

# Key tags
tags = ["100% Local", "Zero Install", "Cross-Platform", "Privacy-First"]
for i, tag in enumerate(tags):
    x = Inches(2.5 + i * 2.2)
    card = add_shape(slide, x, Inches(5.3), Inches(1.8), Inches(0.45), fill_color=None, border_color=ACCENT, border_width=Pt(1.5))
    tf = card.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
    "2026  ·  Enthusiast Guy  ·  Open Source", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "THE PROBLEM", font_size=14, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0))

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
    "How do you showcase branching, loopable music to clients?", font_size=36, color=WHITE, bold=True)

problems = [
    ("🎵", "Branching Audio Is Hard to Demo",
     "Adaptive music with multiple parts, loops, and branch points can't be shown in a simple MP3 player. Clients never hear the full picture."),
    ("💻", "DAWs Are Overkill",
     "Pro Tools, Logic, Ableton — they require installation, steep learning curves, and licenses. Not practical for a quick client audition."),
    ("☁️", "Cloud Tools Risk Privacy",
     "Uploading unreleased music to web platforms raises IP concerns. Composers need local-first solutions that keep audio on the device."),
    ("📦", "No Standard Demo Format",
     "Every composer sends a different folder structure. Clients get confused by dozens of WAV files with no way to hear them in context."),
]

for i, (icon, title, desc) in enumerate(problems):
    y = Inches(2.3 + i * 1.2)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.05), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(0.5), Inches(0.5), icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.6), y + Inches(0.08), Inches(4), Inches(0.35), title, font_size=16, color=ACCENT3, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.45), Inches(10.5), Inches(0.55), desc, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCING S.E.A.M
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "THE SOLUTION", font_size=14, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
    "S.E.A.M — Audio That Speaks For Itself", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
    "A self-contained, browser-based audition player that turns your sample pack folder into an interactive, playable experience — instantly.", font_size=16, color=LIGHT_GRAY)

features_intro = [
    ("🔓", "Open a single HTML file", "No install, no server, no dependencies. Works in Chrome/Edge on any OS.", ACCENT),
    ("📁", "Point to your folder", "Select your sample pack root. S.E.A.M discovers songs, parts, and branches automatically.", ACCENT2),
    ("▶️", "Play, loop & compose", "Audition full songs or individual parts. Drag bricks to build custom sequences. Export results.", ACCENT3),
    ("🔒", "100% Private", "All processing happens locally in your browser. No files ever leave your device.", ACCENT_GREEN),
]

for i, (icon, title, desc, color) in enumerate(features_intro):
    x = Inches(0.5 + i * 3.1)
    y = Inches(3.0)
    card = add_shape(slide, x, y, Inches(2.9), Inches(3.6), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    # Top accent
    add_shape(slide, x + Inches(0.05), y + Inches(0.05), Inches(2.8), Pt(3), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    # Icon
    add_text_box(slide, x, y + Inches(0.3), Inches(2.9), Inches(0.6), icon, font_size=36, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.5), Inches(0.6), title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide, x + Inches(0.2), y + Inches(1.7), Inches(2.5), Inches(1.6), desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "HOW IT WORKS", font_size=14, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0))

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
    "From Folder to Interactive Demo in 3 Clicks", font_size=36, color=WHITE, bold=True)

steps = [
    ("1", "Structure Your Pack",
     "Organize songs in subfolders with the naming convention:\n\n"
     "Song Name 1 → 2, 4.wav  (Part 1, branches to 2 or 4)\n"
     "Song Name 2 → 2, 3.wav  (Part 2, loops or goes to 3)\n"
     "Song Name 3.wav              (Final part, no branch)\n"
     "Song Name.wav                (Optional full mix)",
     ACCENT),
    ("2", "Open & Select",
     "Double-click S.E.A.M's HTML file in any Chromium browser.\n\n"
     "Click 'Select Folder' and point to your sample pack root.\n\n"
     "S.E.A.M auto-discovers all songs, parts, durations,\n"
     "loop points, and branch targets from filenames alone.",
     ACCENT2),
    ("3", "Audition & Export",
     "Play full songs or individual parts with looping.\n\n"
     "Drag bricks to build custom sequences.\n\n"
     "Export audio to WAV, MP3, OGG, or FLAC.\n"
     "Build MP4 demo videos with chapters.\n\n"
     "Session state auto-saves per project folder.",
     ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 4.1)
    y = Inches(2.2)
    # Step number circle
    add_icon_circle(slide, x + Inches(1.35), y, Inches(0.7), color, num, font_size=26)
    # Arrow connector (except last)
    if i < 2:
        add_shape(slide, x + Inches(3.6), y + Inches(0.15), Inches(0.9), Pt(3), fill_color=DARK_GRAY, shape_type=MSO_SHAPE.RECTANGLE)
        add_text_box(slide, x + Inches(3.9), y - Inches(0.05), Inches(0.5), Inches(0.5), "→", font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x, y + Inches(0.9), Inches(3.8), Inches(0.5), title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Card
    card = add_shape(slide, x, y + Inches(1.5), Inches(3.8), Inches(3.5), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(1.65), Inches(3.4), Inches(3.2), desc, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — CORE PLAYBACK FEATURES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "CORE FEATURES", font_size=14, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0))

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Professional Playback, Zero Compromise", font_size=34, color=WHITE, bold=True)

cards = [
    ("▶", "Gapless Playback", "Web Audio API-powered transport with\npre-scheduled segments for seamless\ntransitions between parts.", ACCENT),
    ("🔁", "Smart Looping", "Per-part loop controls auto-detected from\nfilename conventions. Loop once, N times,\nor infinitely with zero gaps.", ACCENT2),
    ("⚡", "Variable Speed", "Rotary knob from -250% to +200% with\n5% snap increments. Reverse playback\nwith pitch-shifted audio.", ACCENT3),
    ("🎚", "Crossfade Engine", "Adjustable 0–20s crossfade between\nplaylist songs with visual knob control.\nSmooth volume curves on transitions.", GOLD),
    ("🎵", "Multi-Track Play", "Play multiple songs simultaneously for\nA/B comparison or layered audition.\nIndependent transport per song.", ACCENT_GREEN),
    ("📊", "Live Waveforms", "Real-time waveform rendering on seek\nbars and bricks. Spectrum analyzer in\nheader. Responsive to audio content.", ACCENT),
]

for i, (icon, title, desc, color) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.1)
    y = Inches(2.1 + row * 2.6)
    add_feature_card(slide, x, y, Inches(3.8), Inches(2.3), icon, title, desc, color)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — BRICK PLAYGROUND & TIMELINE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "POWER FEATURE", font_size=14, color=ACCENT2, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "The Brick Playground — Visual Composition", font_size=34, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(0.5),
    "A unique 2D canvas where you arrange audio segments like building blocks — no other audition tool offers this.", font_size=15, color=LIGHT_GRAY)

# Left panel - features
left_features = [
    "Drag colored 'bricks' representing audio parts onto an infinite 2D canvas",
    "Magnetic snap — bricks auto-cluster when placed within 14px proximity",
    "Break gap system — 36px minimum separation prevents accidental merges",
    "Comb parking lot — all available parts shown on the left, drag to duplicate",
    "Ghost reorder — hold 1s to edit a single brick within a cluster",
    "Cluster operations: merge, unmerge, delete, annotate with title & description",
    "50-step undo/redo stack for non-destructive experimentation",
    "Export any cluster directly to the main timeline for playback",
    "Snow particle idle effect with physics, wind gusts, and deposits on bricks",
    "Live spectrum visualizer with colorful FFT pillars while audio plays",
]

y_start = Inches(2.7)
for i, feat in enumerate(left_features):
    y = y_start + Inches(i * 0.42)
    add_text_box(slide, Inches(1.2), y, Inches(0.3), Inches(0.35), "►", font_size=9, color=ACCENT2)
    add_text_box(slide, Inches(1.5), y, Inches(10.5), Inches(0.4), feat, font_size=12, color=LIGHT_GRAY)

# Bottom highlight box
add_shape(slide, Inches(0.8), Inches(6.95), Inches(11.5), Inches(0.45), fill_color=None, border_color=ACCENT2, border_width=Pt(1.5))
add_text_box(slide, Inches(1.0), Inches(6.97), Inches(11), Inches(0.4),
    "💡  Think of it as a LEGO board for music composition — visual, intuitive, and endlessly reconfigurable",
    font_size=13, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — SEAM PREVIEW & SMART FEATURES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "SMART FEATURES", font_size=14, color=ACCENT3, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), ACCENT3)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Intelligent Audition Tools", font_size=34, color=WHITE, bold=True)

smart_features = [
    ("🔍", "Seam Preview", 
     "A novel audition mode that plays the head and tail of a part while fast-forwarding the middle — perfect for testing loop transition quality without listening to the entire segment. Configurable edge duration (50ms–60s).",
     ACCENT3),
    ("💾", "Session Memory",
     "Every setting is auto-saved per project folder: playlist order, loop counts, sequences, export preferences, brick playground state, encoding options. Reopen any project exactly where you left off via IndexedDB + localStorage.",
     ACCENT),
    ("📡", "Auto-Discovery",
     "Point to any folder and S.E.A.M parses filenames to detect parts, branch targets, loop points, and full mixes automatically. Supports the  Name N → A, B, C.wav  convention with zero configuration needed.",
     ACCENT2),
    ("⚙️", "Encoding Studio",
     "Export to WAV (PCM 16-bit), MP3 (96–320 kbps via lamejs), OGG Vorbis (quality 0–1.0), and FLAC (lossless, level 0–8). Plus built-in 1080p MP4 demo video export with chapter markers.",
     GOLD),
]

for i, (icon, title, desc, color) in enumerate(smart_features):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.0 + row * 2.65)
    card = add_shape(slide, x, y, Inches(6.0), Inches(2.4), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_shape(slide, x + Inches(0.05), y + Inches(0.05), Inches(5.9), Pt(3), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5), icon, font_size=24, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.9), y + Inches(0.25), Inches(4.8), Inches(0.4), title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1.4), desc, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — EXPORT & ENCODING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "EXPORT ENGINE", font_size=14, color=GOLD, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), GOLD)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "From Audition to Deliverable in One Click", font_size=34, color=WHITE, bold=True)

# Audio format cards (2x2)
formats = [
    ("WAV", "Lossless PCM", [
        "16-bit PCM encoding",
        "Source sample rate preserved",
        "Multi-channel support",
        "Instant encoding — no external codec",
    ], ACCENT),
    ("MP3", "Universal Lossy", [
        "Bitrate: 96 – 320 kbps",
        "Sample rate: source, 44.1k, 48k Hz",
        "Channel modes: auto, mono, stereo",
        "Powered by lamejs (lazy-loaded)",
    ], ACCENT3),
    ("OGG", "Open Lossy", [
        "Quality slider: 0.0 – 1.0",
        "Sample rate: source, 44.1k, 48k Hz",
        "Channel modes: auto, mono, stereo",
        "Vorbis encoder (lazy-loaded)",
    ], ACCENT2),
    ("FLAC", "Lossless Compressed", [
        "Compression level: 0 – 8",
        "Sample rate: source, 44.1k, 48k Hz",
        "Channel modes: auto, mono, stereo",
        "libflac WASM runtime (local/offline)",
    ], ACCENT_GREEN),
]

for i, (fmt, sub, features, color) in enumerate(formats):
    col = i % 2
    row = i // 2
    x = Inches(0.7 + col * 6.05)
    y = Inches(2.0 + row * 2.35)
    card = add_shape(slide, x, y, Inches(5.85), Inches(2.15), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_shape(slide, x + Inches(0.05), y + Inches(0.05), Inches(5.75), Pt(4), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.22), Inches(1.4), Inches(0.45), fmt, font_size=26, color=color, bold=True)
    add_text_box(slide, x + Inches(1.7), y + Inches(0.28), Inches(3.9), Inches(0.35), sub, font_size=12, color=MID_GRAY)
    
    for j, feat in enumerate(features):
        fy = y + Inches(0.7 + j * 0.35)
        add_text_box(slide, x + Inches(0.3), fy, Inches(0.3), Inches(0.3), "✓", font_size=13, color=color, bold=True)
        add_text_box(slide, x + Inches(0.6), fy, Inches(5.0), Inches(0.32), feat, font_size=11, color=LIGHT_GRAY)

# Bottom note
add_shape(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(7.02), Inches(11), Inches(0.35),
    "Export full compositions (up to 60 min) or individual parts  •  4 audio formats + MP4 demo video  •  Per-part format memory",
    font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — UI & VISUAL EXPERIENCE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "VISUAL EXPERIENCE", font_size=14, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0))

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Designed to Impress — Not Just to Function", font_size=34, color=WHITE, bold=True)

vis_features = [
    ("🎨", "Ambient Setup Screen",
     "4-layer animated waveform with ribbon effects, particle emission on bass hits, film grain, and cinema mode that dims to darkness when idle. It's a visual experience before you even load audio.",
     ACCENT),
    ("📊", "Live Spectrum Analyzer",
     "420-bar animated FFT spectrum with peak caps in the header. Responds to all active audio in real-time. Color-coded frequency visualization brings audio data to life.",
     ACCENT2),
    ("🎛", "Canvas-Drawn Knobs",
     "Speed and crossfade controlled via realistic rotary knobs with filled arcs, indicator dots, snap detents, click sounds, and touch/scroll/drag support. Blue for forward, red for reverse.",
     ACCENT3),
    ("🎯", "Color-Coded Parts",
     "10-color cycling palette assigns unique colors to each part. Consistent across timeline bricks, part lists, waveforms, and playground. Instant visual identification.",
     GOLD),
    ("❄️", "Snow Particle Physics",
     "Idle detection triggers a snow system with wind gusts, particle physics, deposits on bricks, and melt decay on activity. Up to 23,800 particles with grid-based erosion simulation.",
     ACCENT_GREEN),
    ("📈", "Memory Dashboard",
     "Real-time browser memory monitoring using Performance API. Displays tab memory vs. heap usage with stale-reading detection and 700ms refresh. Know your resource footprint.",
     ACCENT),
]

for i, (icon, title, desc, color) in enumerate(vis_features):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.1)
    y = Inches(2.1 + row * 2.6)
    add_feature_card(slide, x, y, Inches(3.8), Inches(2.3), icon, title, desc, color)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPETITIVE COMPARISON
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "COMPETITIVE LANDSCAPE", font_size=14, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "S.E.A.M vs. The Alternatives", font_size=34, color=WHITE, bold=True)

# Table header
col_w = Inches(10.5) / 5
table_left = Inches(1.4)
header_y = Inches(2.1)
headers = ["Feature", "S.E.A.M", "FMOD Studio", "Wwise", "Audacity"]
header_colors = [MID_GRAY, ACCENT, ACCENT3, ACCENT2, GOLD]

for i, (h, c) in enumerate(zip(headers, header_colors)):
    add_shape(slide, table_left + col_w * i, header_y, col_w, Inches(0.45), fill_color=DARK_GRAY, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, table_left + col_w * i, header_y + Inches(0.02), col_w, Inches(0.4), h, font_size=12, color=c, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ("Zero Install Required", "✓", "✗", "✗", "✗"),
    ("Runs in Browser", "✓", "✗", "✗", "✗"),
    ("Branching Audio Support", "✓", "✓", "✓", "✗"),
    ("Loop-Point Detection", "✓", "Manual", "Manual", "✗"),
    ("Visual Brick Composition", "✓", "✗", "✗", "✗"),
    ("Seam Preview Mode", "✓", "✗", "✗", "✗"),
    ("FLAC Lossless Export", "✓", "✗", "✗", "✗"),
    ("MP4 Demo Video Export", "✓", "✗", "✗", "✗"),
    ("Multi-Format Audio Export", "✓", "✓", "✓", "✓"),
    ("Session Persistence", "✓", "✓", "✓", "✗"),
    ("100% Local / No Upload", "✓", "✓", "✓", "✓"),
    ("Free & Open Source", "✓", "Freemium", "Freemium", "✓"),
    ("Cross-Platform", "✓", "Win/Mac", "Win/Mac", "✓"),
    ("Learning Curve", "None", "Steep", "Steep", "Moderate"),
]

for i, (feat, s, f, w, a) in enumerate(rows):
    y = Inches(2.6 + i * 0.375)
    bg = BG_CARD if i % 2 == 0 else None
    add_comparison_row(slide, table_left, y, Inches(10.5), feat, s, f, w, a, bg)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — USE CASES & TARGET AUDIENCE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "WHO IS IT FOR", font_size=14, color=ACCENT2, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Built for Creators. Loved by Clients.", font_size=34, color=WHITE, bold=True)

personas = [
    ("🎮", "Game Audio Composers",
     "Deliver interactive music demos to game studios. Let clients hear branching cues the way they'll work in-engine — with loops, transitions, and variable playback — without requiring FMOD or Wwise.",
     ACCENT),
    ("🎬", "Film & Media Composers",
     "Showcase adaptive scoring to directors and producers. Demonstrate how musical segments connect and flow without the complexity of a DAW. Simple folder structure, instant playback.",
     ACCENT2),
    ("🏪", "Sample Pack Vendors",
     "Ship interactive demos alongside your sample packs. Customers can audition every part, build custom sequences, and export previews. The ultimate 'try before you buy' experience.",
     ACCENT3),
    ("🎓", "Music Educators",
     "Teach adaptive music composition interactively. Students can visualize branching structures, experiment with arrangements in the brick playground, and understand non-linear music design.",
     GOLD),
]

for i, (icon, title, desc, color) in enumerate(personas):
    x = Inches(0.5 + (i % 2) * 6.3)
    y = Inches(2.1 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(6.0), Inches(2.3), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_shape(slide, x + Inches(0.05), y + Inches(0.05), Inches(5.9), Pt(3), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.5), icon, font_size=28)
    add_text_box(slide, x + Inches(0.85), y + Inches(0.25), Inches(4.9), Inches(0.4), title, font_size=17, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), Inches(5.5), Inches(1.3), desc, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — TECHNICAL ARCHITECTURE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "UNDER THE HOOD", font_size=14, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0))

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Engineered for Performance & Reliability", font_size=34, color=WHITE, bold=True)

tech_items = [
    ("Web Audio API", "Professional audio engine with OfflineAudioContext for seamless stitching, master gain/analyser chain, and gapless pre-scheduled playback.", ACCENT),
    ("File System Access API", "Chromium's modern API for read-only local folder access. No uploads, no servers. IndexedDB stores directory handles for instant project reopening.", ACCENT2),
    ("Canvas 2D Rendering", "All knobs, waveforms, spectrum visualizers, particle systems, and brick playground rendered on HTML5 Canvas with pixel-ratio scaling for Retina displays.", ACCENT3),
    ("Export Pipeline Stack", "On-demand codecs: lamejs (MP3), Vorbis encoder (OGG), libflac WASM (FLAC). MP4 demo export uses WebCodecs + mp4-muxer with chapter metadata injection.", GOLD),
    ("SVG Data URI Caching", "Waveform renders are cached as SVG data URIs per width. ResizeObserver triggers responsive redraws. Column cap at 24,000 for performance.", ACCENT_GREEN),
    ("Session Persistence", "LocalStorage (per folder) saves crossfade, speed, loops, sequences, format prefs (including FLAC), encoding, and playground layout. IndexedDB stores recent project handles.", ACCENT),
]

for i, (title, desc, color) in enumerate(tech_items):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.0 + row * 1.75)
    card = add_shape(slide, x, y, Inches(6.0), Inches(1.55), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_shape(slide, x, y, Pt(4), Inches(1.55), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.12), Inches(5.5), Inches(0.35), title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.55), Inches(5.5), Inches(0.9), desc, font_size=11, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 13 — KEY DIFFERENTIATORS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "WHY S.E.A.M WINS", font_size=14, color=GOLD, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.0), GOLD)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Unique Value Propositions", font_size=34, color=WHITE, bold=True)

differentiators = [
    ("ONLY", "Browser-Native Branching Audio Player",
     "No other tool lets you audition branching, loopable music directly in a browser with zero installation.", ACCENT),
    ("ONLY", "Visual Brick-Based Composition",
     "The Playground's magnetic-snap 2D canvas for spatial audio arrangement is unprecedented in audition tools.", ACCENT2),
    ("ONLY", "Filename-Driven Auto-Discovery",
     "Drop a folder, get a fully parsed interactive player. No metadata files, no configuration, no setup wizards.", ACCENT3),
    ("ONLY", "Seam Preview Audition Mode",
     "Head-FF-Tail loop playback for instant transition quality checking — a workflow innovation for adaptive music.", GOLD),
    ("FREE", "Open Source & Truly Local",
     "No subscription, no account, no telemetry. A single HTML file does everything. Ship it alongside your sample pack.", ACCENT_GREEN),
]

for i, (badge, title, desc, color) in enumerate(differentiators):
    y = Inches(2.0 + i * 1.05)
    card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.9), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    # Badge
    badge_shape = add_shape(slide, Inches(1.0), y + Inches(0.15), Inches(0.75), Inches(0.5), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = badge_shape.text_frame
    p = tf.paragraphs[0]
    p.text = badge
    p.font.size = Pt(10)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text_box(slide, Inches(2.0), y + Inches(0.1), Inches(10), Inches(0.35), title, font_size=16, color=WHITE, bold=True)
    # Desc
    add_text_box(slide, Inches(2.0), y + Inches(0.45), Inches(10), Inches(0.4), desc, font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 14 — MP4 DEMO VIDEO EXPORT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
    "SHOWCASE TOOLING", font_size=14, color=ACCENT2, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(2.4), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
    "Built-In 1080p MP4 Demo Video Export", font_size=34, color=WHITE, bold=True)

video_items = [
    ("🎬", "Offline Render Pipeline",
     "Canvas frames + WebCodecs VideoEncoder (H.264) and AudioEncoder (AAC) are muxed locally into MP4. No upload, no cloud render queue.",
     ACCENT2),
    ("🏷", "Chapters + Description",
     "Injects Nero chapter markers into MP4 and generates a YouTube-style description.txt with track timeline markers for posting and review.",
     ACCENT_GREEN),
    ("🎨", "Branding Controls",
     "Choose bundled OFL fonts, per-zone text size tiers, corner credit (optional year), and custom image/video backgrounds with crossfaded loops.",
     ACCENT3),
    ("🧠", "Smart Visual Adaptation",
     "Preview mode and Auto palette derive contrast-safe UI colors from background images or live video frames, while preserving a fixed Basic theme option.",
     GOLD),
]

for i, (icon, title, desc, color) in enumerate(video_items):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.0 + row * 2.4)
    card = add_shape(slide, x, y, Inches(6.0), Inches(2.15), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
    add_shape(slide, x + Inches(0.05), y + Inches(0.05), Inches(5.9), Pt(3), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.45), icon, font_size=24, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.85), y + Inches(0.22), Inches(4.9), Inches(0.35), title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.72), Inches(5.5), Inches(1.25), desc, font_size=11, color=LIGHT_GRAY)

add_shape(slide, Inches(0.8), Inches(6.95), Inches(11.5), Inches(0.4), fill_color=BG_CARD, border_color=DARK_GRAY, border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(6.98), Inches(11.0), Inches(0.35),
    "Default output profile: 1920×1080, up to 24 FPS, H.264 video + AAC audio, with embedded chapter metadata.",
    font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING / CTA
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative gradient bars (same as title)
for i, c in enumerate([ACCENT2, ACCENT, ACCENT_GREEN, GOLD]):
    add_shape(slide, Inches(1.5 + i * 2.7), Inches(0.15), Inches(2.2), Pt(4), fill_color=c, shape_type=MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
    "S.E.A.M", font_size=64, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER, font_name='Segoe UI Light')

add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.6),
    "Your Music. Their Experience. Zero Friction.", font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.0), Inches(3.4), Inches(3.333), ACCENT)

# Key stats
stats = [
    ("0", "Dependencies", ACCENT),
    ("< 1s", "Load Time", ACCENT2),
    ("4 + MP4", "Export Formats", ACCENT3),
    ("100%", "Privacy", ACCENT_GREEN),
    ("∞", "Creativity", GOLD),
]

for i, (val, label, color) in enumerate(stats):
    x = Inches(1.2 + i * 2.2)
    y = Inches(3.9)
    add_text_box(slide, x, y, Inches(2.0), Inches(0.7), val, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.65), Inches(2.0), Inches(0.35), label, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# CTA buttons
cta_items = [
    ("Try S.E.A.M Today — Just Open index.html", ACCENT),
    ("github.com/EnthusiastGuy/S.E.A.M-Audio-Audition", ACCENT2),
]

for i, (text, color) in enumerate(cta_items):
    y = Inches(5.5 + i * 0.65)
    card = add_shape(slide, Inches(3.0), y, Inches(7.333), Inches(0.5), fill_color=None, border_color=color, border_width=Pt(2))
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10), Inches(0.5),
    "2026  ·  Enthusiast Guy  ·  Proprietary License, free to use", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom decorative bars
for i, c in enumerate([ACCENT2, ACCENT, ACCENT_GREEN, GOLD]):
    add_shape(slide, Inches(1.5 + i * 2.7), Inches(7.3), Inches(2.2), Pt(4), fill_color=c, shape_type=MSO_SHAPE.RECTANGLE)

# ── SAVE ──────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'SEAM_Audio_Audition_Presentation.pptx')
prs.save(output_path)
print(f"✓ Presentation saved to: {output_path}")
print(f"  Slides: {len(prs.slides)}")
